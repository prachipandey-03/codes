"""
Time Management - Professional PowerPoint Presentation Generator
Modern, Futuristic, Minimal Design with Soft Pastel Colors
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
import math
import os

# ============================================================
# COLOR PALETTE - Soft Pastels
# ============================================================
class Colors:
    WHITE = RGBColor(255, 255, 255)
    SOFT_WHITE = RGBColor(248, 249, 252)
    LIGHT_BG = RGBColor(245, 247, 255)
    
    # Primary pastels
    SOFT_BLUE = RGBColor(130, 170, 255)
    SOFT_PURPLE = RGBColor(167, 139, 250)
    SOFT_LAVENDER = RGBColor(196, 181, 253)
    SOFT_PINK = RGBColor(244, 163, 200)
    SOFT_TEAL = RGBColor(110, 210, 210)
    SOFT_MINT = RGBColor(134, 239, 200)
    SOFT_PEACH = RGBColor(251, 191, 146)
    SOFT_CORAL = RGBColor(248, 150, 150)
    SOFT_SKY = RGBColor(147, 197, 253)
    SOFT_INDIGO = RGBColor(129, 140, 248)
    
    # Text colors
    DARK_TEXT = RGBColor(45, 50, 80)
    MEDIUM_TEXT = RGBColor(80, 90, 120)
    LIGHT_TEXT = RGBColor(120, 130, 160)
    
    # Accent
    ACCENT_BLUE = RGBColor(99, 130, 255)
    ACCENT_PURPLE = RGBColor(139, 92, 246)
    LIGHT_GRAY = RGBColor(226, 232, 240)
    VERY_LIGHT_PURPLE = RGBColor(237, 233, 254)
    VERY_LIGHT_BLUE = RGBColor(219, 234, 254)
    VERY_LIGHT_PINK = RGBColor(252, 231, 243)
    VERY_LIGHT_MINT = RGBColor(209, 250, 229)
    VERY_LIGHT_PEACH = RGBColor(254, 243, 199)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ============================================================
# HELPER FUNCTIONS
# ============================================================

def add_background(slide, color=Colors.WHITE):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, shape_type, left, top, width, height, fill_color=None, line_color=None, line_width=Pt(0)):
    """Add a shape with optional fill and line."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color, line_color=None, shadow=False):
    """Add a rounded rectangle."""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, fill_color, line_color)
    if shadow:
        shape.shadow.inherit = False
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, font_color=Colors.DARK_TEXT, 
                 bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_points(slide, left, top, width, height, items, font_size=16, 
                      font_color=Colors.DARK_TEXT, icon="●", spacing=Pt(12), font_name="Calibri"):
    """Add bullet points with custom icon."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{icon}  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.name = font_name
        p.space_after = spacing
        p.space_before = Pt(4)
    return txBox

def add_circle_icon(slide, cx, cy, radius, fill_color, icon_text="", text_color=Colors.WHITE, font_size=20):
    """Add a circle with centered text as an icon."""
    shape = add_shape(slide, MSO_SHAPE.OVAL, 
                      cx - radius, cy - radius, radius * 2, radius * 2, fill_color)
    if icon_text:
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = icon_text
        p.font.size = Pt(font_size)
        p.font.color.rgb = text_color
        p.font.bold = True
        p.font.name = "Segoe UI Emoji"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    return shape

def add_decorative_dots(slide, x, y, rows, cols, spacing, color, size=Inches(0.08)):
    """Add decorative dot pattern."""
    for r in range(rows):
        for c in range(cols):
            add_shape(slide, MSO_SHAPE.OVAL,
                     x + c * spacing, y + r * spacing, size, size, color)

def add_gradient_bar(slide, left, top, width, height, color):
    """Add a gradient-style bar."""
    shape = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, color)
    return shape

def add_slide_number(slide, num, total=11):
    """Add slide number at bottom right."""
    add_text_box(slide, Inches(11.8), Inches(6.9), Inches(1.2), Inches(0.4),
                f"{num} / {total}", font_size=10, font_color=Colors.LIGHT_TEXT, 
                alignment=PP_ALIGN.RIGHT)

def add_section_header(slide, title, subtitle="", y_offset=Inches(0.4)):
    """Add a section header with accent line."""
    # Accent bar
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 
              Inches(0.6), y_offset, Inches(0.08), Inches(0.6), Colors.ACCENT_PURPLE)
    
    # Title
    add_text_box(slide, Inches(0.9), y_offset - Inches(0.05), Inches(8), Inches(0.6),
                title, font_size=30, font_color=Colors.DARK_TEXT, bold=True)
    
    if subtitle:
        add_text_box(slide, Inches(0.9), y_offset + Inches(0.55), Inches(8), Inches(0.4),
                    subtitle, font_size=14, font_color=Colors.LIGHT_TEXT)

def add_top_accent_strip(slide):
    """Add a thin gradient-like accent strip at top."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.06), Colors.SOFT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.06), SLIDE_WIDTH, Inches(0.03), Colors.SOFT_PURPLE)

def add_card(slide, left, top, width, height, fill_color=Colors.WHITE, border_color=Colors.LIGHT_GRAY):
    """Add a card-style container."""
    shape = add_rounded_rect(slide, left, top, width, height, fill_color, border_color)
    shape.line.width = Pt(1)
    return shape


# ============================================================
# SLIDE 1: TITLE SLIDE
# ============================================================
def create_slide_1(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_background(slide, Colors.WHITE)
    
    # Top accent strip
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), Colors.SOFT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.08), SLIDE_WIDTH, Inches(0.04), Colors.SOFT_LAVENDER)
    
    # Decorative circles (top-right) - kept inside page
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10.3), Inches(0.2), Inches(2.5), Inches(2.5), Colors.VERY_LIGHT_BLUE)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(11.0), Inches(0.5), Inches(1.5), Inches(1.5), Colors.VERY_LIGHT_PURPLE)
    
    # Decorative circles (bottom-left) - kept inside page
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.2), Inches(5.0), Inches(2.2), Inches(2.2), Colors.VERY_LIGHT_PINK)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), Inches(5.5), Inches(1.5), Inches(1.5), Colors.VERY_LIGHT_MINT)
    
    # Decorative dots
    add_decorative_dots(slide, Inches(10), Inches(5.5), 4, 4, Inches(0.25), Colors.SOFT_LAVENDER)
    
    # Clock icon (large circle)
    add_circle_icon(slide, Inches(10.5), Inches(3.3), Inches(0.6), Colors.SOFT_BLUE, "🕐", Colors.WHITE, 28)
    
    # Calendar icon
    add_circle_icon(slide, Inches(11.6), Inches(4.2), Inches(0.5), Colors.SOFT_PURPLE, "📅", Colors.WHITE, 24)
    
    # Hourglass icon
    add_circle_icon(slide, Inches(10.0), Inches(4.7), Inches(0.45), Colors.SOFT_PINK, "⏳", Colors.WHITE, 22)
    
    # Main Title
    add_text_box(slide, Inches(0.8), Inches(1.2), Inches(9), Inches(1.2),
                "Time Management", font_size=52, font_color=Colors.DARK_TEXT, bold=True)
    
    # Subtitle
    add_text_box(slide, Inches(0.8), Inches(2.3), Inches(9), Inches(0.6),
                "Work Smart, Save Time", font_size=24, font_color=Colors.ACCENT_PURPLE, bold=False)
    
    # Divider line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(3.1), Inches(3), Inches(0.04), Colors.SOFT_BLUE)
    
    # Team heading
    add_text_box(slide, Inches(0.8), Inches(3.4), Inches(5), Inches(0.4),
                "Team Members", font_size=16, font_color=Colors.ACCENT_BLUE, bold=True)
    
    # Team members - Two columns
    team_left = [
        "Prachi Pandey – S25CSEU1097",
        "Nandini Arora – S25CSEU0631",
        "Naitik Jain – S25CSEU0438",
        "Tushar Joshi – S25CSEU2000",
    ]
    team_right = [
        "Kushagra Tailor – S25CSEU1720",
        "Gaurav Dubey – A25ARIU0059",
        "Abhinav Bhatnagar – M25BBAU0065",
        "Arjun Singh – S25CSEU2221",
    ]
    
    y_start = Inches(3.9)
    for i, member in enumerate(team_left):
        # Small circle bullet
        add_shape(slide, MSO_SHAPE.OVAL, Inches(0.9), y_start + Inches(i * 0.4) + Inches(0.06), 
                  Inches(0.12), Inches(0.12), Colors.SOFT_BLUE)
        add_text_box(slide, Inches(1.15), y_start + Inches(i * 0.4), Inches(4), Inches(0.35),
                    member, font_size=13, font_color=Colors.MEDIUM_TEXT)
    
    for i, member in enumerate(team_right):
        add_shape(slide, MSO_SHAPE.OVAL, Inches(5.5), y_start + Inches(i * 0.4) + Inches(0.06),
                  Inches(0.12), Inches(0.12), Colors.SOFT_PURPLE)
        add_text_box(slide, Inches(5.75), y_start + Inches(i * 0.4), Inches(4.5), Inches(0.35),
                    member, font_size=13, font_color=Colors.MEDIUM_TEXT)
    
    # Bottom accent
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12), Colors.SOFT_BLUE)
    
    add_slide_number(slide, 1)


# ============================================================
# SLIDE 2: INTRODUCTION
# ============================================================
def create_slide_2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    # Decorative circle - kept inside page
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10.8), Inches(5.0), Inches(2.2), Inches(2.2), Colors.VERY_LIGHT_BLUE)
    
    add_section_header(slide, "Introduction", "Understanding the basics of time management")
    
    # Content cards - left side
    items = [
        ("🕐", "Time management means\nusing time wisely", Colors.VERY_LIGHT_BLUE, Colors.SOFT_BLUE),
        ("✅", "Helps us finish work\non time", Colors.VERY_LIGHT_MINT, Colors.SOFT_MINT),
        ("🎓", "Very important for\nstudents & professionals", Colors.VERY_LIGHT_PURPLE, Colors.SOFT_PURPLE),
        ("📈", "Leads to better\nresults and growth", Colors.VERY_LIGHT_PEACH, Colors.SOFT_PEACH),
    ]
    
    for i, (icon, text, bg_color, accent) in enumerate(items):
        card_x = Inches(0.7) + Inches(i * 2.9)
        card_y = Inches(1.8)
        
        # Card background
        card = add_card(slide, card_x, card_y, Inches(2.6), Inches(2.2), bg_color)
        card.line.fill.background()
        
        # Top accent
        add_shape(slide, MSO_SHAPE.RECTANGLE, card_x, card_y, Inches(2.6), Inches(0.06), accent)
        
        # Icon circle
        add_circle_icon(slide, card_x + Inches(1.3), card_y + Inches(0.7), Inches(0.35), accent, icon, Colors.WHITE, 22)
        
        # Text
        add_text_box(slide, card_x + Inches(0.2), card_y + Inches(1.2), Inches(2.2), Inches(0.9),
                    text, font_size=14, font_color=Colors.DARK_TEXT, alignment=PP_ALIGN.CENTER)
    
    # Timeline at bottom
    add_text_box(slide, Inches(0.7), Inches(4.4), Inches(5), Inches(0.4),
                "⏱ Simple Timeline: How Time Management Works", font_size=15, 
                font_color=Colors.ACCENT_PURPLE, bold=True)
    
    # Timeline line
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.45), Inches(10), Inches(0.04), Colors.SOFT_BLUE)
    
    # Timeline nodes
    timeline_items = [
        ("Plan", "📋", Colors.SOFT_BLUE),
        ("Prioritize", "🎯", Colors.SOFT_PURPLE),
        ("Execute", "⚡", Colors.SOFT_PINK),
        ("Review", "🔍", Colors.SOFT_TEAL),
        ("Improve", "🚀", Colors.SOFT_MINT),
    ]
    
    for i, (label, icon, color) in enumerate(timeline_items):
        x = Inches(1.5) + Inches(i * 2.5)
        # Node circle
        add_circle_icon(slide, x + Inches(0.25), Inches(5.45), Inches(0.25), color, icon, Colors.WHITE, 14)
        # Label
        add_text_box(slide, x - Inches(0.2), Inches(5.85), Inches(1.0), Inches(0.35),
                    label, font_size=12, font_color=Colors.DARK_TEXT, alignment=PP_ALIGN.CENTER, bold=True)
    
    add_slide_number(slide, 2)


# ============================================================
# SLIDE 3: IMPORTANCE - with Bar Chart
# ============================================================
def create_slide_3(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Importance of Time Management", "Why it matters for your success")
    
    # Left side: Bullet points with icons
    importance_items = [
        ("⏰", "Helps Save Time", "Finish tasks faster with proper planning", Colors.SOFT_BLUE),
        ("😌", "Reduces Stress", "Less worry about deadlines", Colors.SOFT_MINT),
        ("🎯", "Improves Focus", "Stay on track with your goals", Colors.SOFT_PURPLE),
        ("🏆", "Achieve Goals", "Reach your targets step by step", Colors.SOFT_PEACH),
    ]
    
    for i, (icon, title, desc, color) in enumerate(importance_items):
        y = Inches(1.7) + Inches(i * 1.25)
        
        # Icon circle
        add_circle_icon(slide, Inches(1.2), y + Inches(0.3), Inches(0.32), color, icon, Colors.WHITE, 18)
        
        # Title
        add_text_box(slide, Inches(1.8), y, Inches(4), Inches(0.35),
                    title, font_size=17, font_color=Colors.DARK_TEXT, bold=True)
        # Description
        add_text_box(slide, Inches(1.8), y + Inches(0.35), Inches(4), Inches(0.3),
                    desc, font_size=12, font_color=Colors.LIGHT_TEXT)
    
    # Right side: Bar Chart - Productivity Increase
    # Chart background card
    add_card(slide, Inches(6.5), Inches(1.5), Inches(6.2), Inches(5.2), Colors.SOFT_WHITE, Colors.LIGHT_GRAY)
    
    add_text_box(slide, Inches(6.8), Inches(1.7), Inches(5), Inches(0.4),
                "📊 Productivity Increase", font_size=15, font_color=Colors.DARK_TEXT, bold=True)
    
    # Bar chart using shapes
    categories = ["Without\nPlanning", "Basic\nPlanning", "To-Do\nList", "Full Time\nMgmt"]
    values = [35, 55, 75, 95]
    bar_colors = [Colors.SOFT_CORAL, Colors.SOFT_PEACH, Colors.SOFT_SKY, Colors.SOFT_BLUE]
    
    chart_base_y = Inches(5.9)
    chart_base_x = Inches(7.2)
    max_height = Inches(3.0)
    bar_width = Inches(1.0)
    bar_spacing = Inches(1.25)
    
    for i, (cat, val, color) in enumerate(zip(categories, values, bar_colors)):
        bar_height = max_height * (val / 100)
        x = chart_base_x + Inches(i * 1.25)
        y = chart_base_y - bar_height
        
        # Bar
        bar = add_rounded_rect(slide, x, y, bar_width, bar_height, color)
        bar.line.fill.background()
        
        # Value label
        add_text_box(slide, x, y - Inches(0.3), bar_width, Inches(0.3),
                    f"{val}%", font_size=13, font_color=Colors.DARK_TEXT, alignment=PP_ALIGN.CENTER, bold=True)
        
        # Category label
        add_text_box(slide, x - Inches(0.1), chart_base_y + Inches(0.1), Inches(1.2), Inches(0.5),
                    cat, font_size=10, font_color=Colors.MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
    
    # Baseline
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(7.0), chart_base_y, Inches(5.3), Inches(0.02), Colors.LIGHT_GRAY)
    
    add_slide_number(slide, 3)


# ============================================================
# SLIDE 4: COMMON PROBLEMS
# ============================================================
def create_slide_4(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Common Problems", "Challenges that waste our time")
    
    # Problem cards in 2x2 grid
    problems = [
        ("😴", "Procrastination", "Delaying work\nuntil last moment", Colors.SOFT_CORAL, Colors.VERY_LIGHT_PINK),
        ("📱", "Too Much Phone", "Social media and\nendless scrolling", Colors.SOFT_BLUE, Colors.VERY_LIGHT_BLUE),
        ("📋", "No Planning", "Starting work without\na clear plan", Colors.SOFT_PEACH, Colors.VERY_LIGHT_PEACH),
        ("❓", "Priority Confusion", "Not knowing what\nto do first", Colors.SOFT_PURPLE, Colors.VERY_LIGHT_PURPLE),
    ]
    
    positions = [
        (Inches(0.7), Inches(1.7)),
        (Inches(6.6), Inches(1.7)),
        (Inches(0.7), Inches(4.3)),
        (Inches(6.6), Inches(4.3)),
    ]
    
    for (icon, title, desc, accent, bg), (px, py) in zip(problems, positions):
        card_w = Inches(5.8)
        card_h = Inches(2.3)
        
        # Card
        card = add_card(slide, px, py, card_w, card_h, bg)
        card.line.fill.background()
        
        # Left accent bar
        add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, px, py, Inches(0.07), card_h, accent)
        
        # Icon
        add_circle_icon(slide, px + Inches(1.0), py + Inches(1.0), Inches(0.45), accent, icon, Colors.WHITE, 26)
        
        # Title
        add_text_box(slide, px + Inches(1.7), py + Inches(0.4), Inches(3.5), Inches(0.4),
                    title, font_size=20, font_color=Colors.DARK_TEXT, bold=True)
        
        # Description
        add_text_box(slide, px + Inches(1.7), py + Inches(0.9), Inches(3.8), Inches(0.9),
                    desc, font_size=14, font_color=Colors.MEDIUM_TEXT)
        
        # Warning indicator
        add_shape(slide, MSO_SHAPE.OVAL, px + card_w - Inches(0.6), py + Inches(0.3), 
                  Inches(0.25), Inches(0.25), accent)
    
    # Center divider decoration
    add_shape(slide, MSO_SHAPE.OVAL, Inches(6.15), Inches(3.2), Inches(0.6), Inches(0.6), Colors.VERY_LIGHT_PURPLE)
    add_text_box(slide, Inches(6.15), Inches(3.25), Inches(0.6), Inches(0.5),
                "⚠️", font_size=18, alignment=PP_ALIGN.CENTER)
    
    add_slide_number(slide, 4)


# ============================================================
# SLIDE 5: KEY TECHNIQUES - Flow/Checklist
# ============================================================
def create_slide_5(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Key Techniques", "Simple steps to manage your time better")
    
    # Flow diagram - horizontal steps
    steps = [
        ("1", "Make a\nTo-Do List", "📝", "Write down all\nyour tasks", Colors.SOFT_BLUE),
        ("2", "Set\nPriorities", "🎯", "Decide what is\nmost important", Colors.SOFT_PURPLE),
        ("3", "Divide Time\nfor Tasks", "⏰", "Give each task\na time slot", Colors.SOFT_PINK),
        ("4", "Set\nDeadlines", "📅", "Add a due date\nfor each task", Colors.SOFT_TEAL),
        ("5", "Review &\nImprove", "✅", "Check progress\nand adjust", Colors.SOFT_MINT),
    ]
    
    start_x = Inches(0.5)
    card_width = Inches(2.2)
    spacing = Inches(0.35)
    y_top = Inches(1.8)
    
    for i, (num, title, icon, desc, color) in enumerate(steps):
        x = start_x + Inches(i * 2.55)
        
        # Card
        card = add_card(slide, x, y_top, card_width, Inches(3.8), Colors.WHITE, Colors.LIGHT_GRAY)
        
        # Top colored strip
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y_top, card_width, Inches(0.06), color)
        
        # Step number circle
        add_circle_icon(slide, x + card_width/2, y_top + Inches(0.6), Inches(0.3), color, num, Colors.WHITE, 18)
        
        # Icon
        add_text_box(slide, x, y_top + Inches(1.1), card_width, Inches(0.5),
                    icon, font_size=30, alignment=PP_ALIGN.CENTER)
        
        # Title
        add_text_box(slide, x + Inches(0.15), y_top + Inches(1.7), card_width - Inches(0.3), Inches(0.7),
                    title, font_size=15, font_color=Colors.DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Description
        add_text_box(slide, x + Inches(0.15), y_top + Inches(2.5), card_width - Inches(0.3), Inches(0.8),
                    desc, font_size=11, font_color=Colors.MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
        
        # Arrow between cards
        if i < len(steps) - 1:
            arrow_x = x + card_width + Inches(0.05)
            add_text_box(slide, arrow_x, y_top + Inches(1.5), Inches(0.3), Inches(0.4),
                        "→", font_size=22, font_color=Colors.SOFT_BLUE, alignment=PP_ALIGN.CENTER, bold=True)
    
    # Bottom checklist summary
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), Colors.VERY_LIGHT_BLUE)
    checklist = "✅ List Tasks    →    ✅ Prioritize    →    ✅ Allocate Time    →    ✅ Set Deadlines    →    ✅ Review Progress"
    add_text_box(slide, Inches(0.8), Inches(6.2), Inches(12), Inches(0.5),
                checklist, font_size=14, font_color=Colors.ACCENT_BLUE, alignment=PP_ALIGN.CENTER, bold=True)
    
    add_slide_number(slide, 5)


# ============================================================
# SLIDE 6: POPULAR METHODS - with Eisenhower Matrix
# ============================================================
def create_slide_6(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Popular Methods", "Proven techniques used by successful people")
    
    # Left side: Method cards
    methods = [
        ("🍅", "Pomodoro Technique", "Work for 25 min + 5 min break\nStay focused in short bursts", Colors.SOFT_CORAL),
        ("📊", "80/20 Rule (Pareto)", "Focus on 20% tasks that give\n80% of results", Colors.SOFT_BLUE),
    ]
    
    for i, (icon, title, desc, color) in enumerate(methods):
        y = Inches(1.7) + Inches(i * 2.3)
        
        # Card background
        card = add_card(slide, Inches(0.5), y, Inches(5.0), Inches(2.0), Colors.WHITE, Colors.LIGHT_GRAY)
        
        # Accent
        add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.06), Inches(2.0), color)
        
        # Icon
        add_circle_icon(slide, Inches(1.3), y + Inches(0.9), Inches(0.4), color, icon, Colors.WHITE, 22)
        
        # Title
        add_text_box(slide, Inches(2.0), y + Inches(0.3), Inches(3.2), Inches(0.4),
                    title, font_size=18, font_color=Colors.DARK_TEXT, bold=True)
        
        # Desc
        add_text_box(slide, Inches(2.0), y + Inches(0.8), Inches(3.2), Inches(0.9),
                    desc, font_size=12, font_color=Colors.MEDIUM_TEXT)
    
    # Right side: Eisenhower Matrix
    add_text_box(slide, Inches(6.2), Inches(1.5), Inches(6), Inches(0.4),
                "📋 Eisenhower Matrix", font_size=17, font_color=Colors.DARK_TEXT, bold=True)
    
    matrix_x = Inches(6.2)
    matrix_y = Inches(2.1)
    cell_w = Inches(3.3)
    cell_h = Inches(2.3)
    
    matrix_items = [
        ("DO FIRST", "Urgent + Important\n\n• Crisis tasks\n• Deadlines today", Colors.SOFT_CORAL, Colors.VERY_LIGHT_PINK),
        ("SCHEDULE", "Not Urgent + Important\n\n• Study planning\n• Skill building", Colors.SOFT_BLUE, Colors.VERY_LIGHT_BLUE),
        ("DELEGATE", "Urgent + Not Important\n\n• Some emails\n• Phone calls", Colors.SOFT_PEACH, Colors.VERY_LIGHT_PEACH),
        ("ELIMINATE", "Not Urgent + Not Important\n\n• Social media\n• Time wasters", Colors.SOFT_PURPLE, Colors.VERY_LIGHT_PURPLE),
    ]
    
    positions_matrix = [
        (matrix_x, matrix_y),
        (matrix_x + cell_w + Inches(0.15), matrix_y),
        (matrix_x, matrix_y + cell_h + Inches(0.15)),
        (matrix_x + cell_w + Inches(0.15), matrix_y + cell_h + Inches(0.15)),
    ]
    
    for (title, content, accent, bg), (mx, my) in zip(matrix_items, positions_matrix):
        # Cell
        cell = add_card(slide, mx, my, cell_w, cell_h, bg)
        cell.line.fill.background()
        
        # Top bar
        add_shape(slide, MSO_SHAPE.RECTANGLE, mx, my, cell_w, Inches(0.05), accent)
        
        # Title
        add_text_box(slide, mx + Inches(0.15), my + Inches(0.15), cell_w - Inches(0.3), Inches(0.35),
                    title, font_size=13, font_color=accent, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Content
        add_text_box(slide, mx + Inches(0.15), my + Inches(0.5), cell_w - Inches(0.3), Inches(1.6),
                    content, font_size=10, font_color=Colors.MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
    
    # Axis labels
    add_text_box(slide, Inches(6.5), Inches(1.85), Inches(3), Inches(0.25),
                "← URGENT", font_size=9, font_color=Colors.LIGHT_TEXT, bold=True)
    add_text_box(slide, Inches(9.8), Inches(1.85), Inches(3), Inches(0.25),
                "NOT URGENT →", font_size=9, font_color=Colors.LIGHT_TEXT, bold=True)
    
    add_slide_number(slide, 6)


# ============================================================
# SLIDE 7: TOOLS FOR TIME MANAGEMENT
# ============================================================
def create_slide_7(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Tools for Time Management", "Digital tools that help you stay organized")
    
    # Tool cards
    tools = [
        ("📅", "Google\nCalendar", "Schedule events\n& set reminders", Colors.SOFT_BLUE, Colors.VERY_LIGHT_BLUE,
         ["Free to use", "Syncs all devices", "Color coding", "Reminders"]),
        ("📝", "Notion", "Notes, tasks &\nproject planning", Colors.SOFT_PURPLE, Colors.VERY_LIGHT_PURPLE,
         ["All-in-one tool", "Templates", "Team sharing", "Databases"]),
        ("✅", "To-Do\nApps", "Track daily tasks\n& checklists", Colors.SOFT_MINT, Colors.VERY_LIGHT_MINT,
         ["Simple lists", "Due dates", "Priorities", "Daily planner"]),
        ("⏱", "Pomodoro\nTimers", "Focus timer with\nwork + breaks", Colors.SOFT_PEACH, Colors.VERY_LIGHT_PEACH,
         ["25 min focus", "5 min breaks", "Track sessions", "Stay focused"]),
    ]
    
    for i, (icon, name, desc, accent, bg, features) in enumerate(tools):
        x = Inches(0.5) + Inches(i * 3.2)
        y = Inches(1.7)
        card_w = Inches(2.9)
        
        # Card
        card = add_card(slide, x, y, card_w, Inches(5.0), bg)
        card.line.fill.background()
        
        # Top accent
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.06), accent)
        
        # Icon
        add_circle_icon(slide, x + card_w/2, y + Inches(0.7), Inches(0.45), accent, icon, Colors.WHITE, 26)
        
        # Name
        add_text_box(slide, x + Inches(0.1), y + Inches(1.3), card_w - Inches(0.2), Inches(0.6),
                    name, font_size=17, font_color=Colors.DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
        
        # Description
        add_text_box(slide, x + Inches(0.1), y + Inches(1.9), card_w - Inches(0.2), Inches(0.6),
                    desc, font_size=11, font_color=Colors.MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)
        
        # Divider
        add_shape(slide, MSO_SHAPE.RECTANGLE, x + Inches(0.3), y + Inches(2.6), card_w - Inches(0.6), Inches(0.02), Colors.LIGHT_GRAY)
        
        # Features
        for j, feat in enumerate(features):
            fy = y + Inches(2.8) + Inches(j * 0.4)
            add_shape(slide, MSO_SHAPE.OVAL, x + Inches(0.3), fy + Inches(0.06), 
                      Inches(0.1), Inches(0.1), accent)
            add_text_box(slide, x + Inches(0.55), fy, card_w - Inches(0.8), Inches(0.3),
                        feat, font_size=11, font_color=Colors.MEDIUM_TEXT)
    
    add_slide_number(slide, 7)


# ============================================================
# SLIDE 8: TIPS FOR STUDENTS
# ============================================================
def create_slide_8(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Tips for Students", "Easy habits to manage your time better")
    
    # Decorative element - kept inside page
    add_shape(slide, MSO_SHAPE.OVAL, Inches(10.8), Inches(0.2), Inches(2.2), Inches(2.2), Colors.VERY_LIGHT_PURPLE)
    
    # Checklist style cards
    tips = [
        ("1", "Do One Task at a Time", "Multitasking reduces quality.\nFocus on a single task for best results.", "🎯", Colors.SOFT_BLUE),
        ("2", "Take Small Breaks", "Short breaks help recharge your brain.\nTry 5-10 minutes between study sessions.", "☕", Colors.SOFT_MINT),
        ("3", "Plan Your Day", "Spend 5 minutes each morning\nwriting your plan for the day.", "📋", Colors.SOFT_PURPLE),
        ("4", "Avoid Distractions", "Keep phone away while studying.\nFind a quiet place to focus.", "🔕", Colors.SOFT_CORAL),
        ("5", "Set Small Goals", "Break big tasks into small steps.\nCelebrate each small win!", "🏆", Colors.SOFT_PEACH),
    ]
    
    for i, (num, title, desc, icon, color) in enumerate(tips):
        y = Inches(1.5) + Inches(i * 1.15)
        
        # Row background (alternating subtle)
        if i % 2 == 0:
            add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(1.0), 
                     Colors.SOFT_WHITE)
        
        # Checkbox style
        checkbox = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), y + Inches(0.2), 
                            Inches(0.45), Inches(0.45), color)
        checkbox.line.fill.background()
        add_text_box(slide, Inches(0.7), y + Inches(0.18), Inches(0.45), Inches(0.45),
                    "✓", font_size=18, font_color=Colors.WHITE, alignment=PP_ALIGN.CENTER, bold=True)
        
        # Title
        add_text_box(slide, Inches(1.4), y + Inches(0.1), Inches(4), Inches(0.35),
                    title, font_size=17, font_color=Colors.DARK_TEXT, bold=True)
        
        # Description
        add_text_box(slide, Inches(1.4), y + Inches(0.5), Inches(6), Inches(0.5),
                    desc, font_size=11, font_color=Colors.MEDIUM_TEXT)
        
        # Right icon - kept inside page
        add_circle_icon(slide, Inches(12.0), y + Inches(0.45), Inches(0.3), color, icon, Colors.WHITE, 18)
    
    add_slide_number(slide, 8)


# ============================================================
# SLIDE 9: REAL-LIFE EXAMPLE - Before vs After
# ============================================================
def create_slide_9(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Real-Life Example", "A student's journey with time management")
    
    # BEFORE card
    before_x = Inches(0.5)
    before_y = Inches(1.6)
    card_w = Inches(5.8)
    card_h = Inches(5.2)
    
    card = add_card(slide, before_x, before_y, card_w, card_h, Colors.VERY_LIGHT_PINK)
    card.line.fill.background()
    
    # Before header
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, before_x, before_y, card_w, Inches(0.7), Colors.SOFT_CORAL)
    add_text_box(slide, before_x + Inches(0.3), before_y + Inches(0.12), card_w - Inches(0.6), Inches(0.5),
                "😰 BEFORE (No Time Management)", font_size=17, font_color=Colors.WHITE, bold=True)
    
    before_items = [
        ("❌", "Wakes up late, no plan for the day"),
        ("❌", "Scrolls phone for 2 hours"),
        ("❌", "Starts homework at night"),
        ("❌", "Feels stressed and tired"),
        ("❌", "Misses deadlines often"),
        ("❌", "No time for fun or hobbies"),
    ]
    
    for i, (icon, text) in enumerate(before_items):
        y = before_y + Inches(1.0) + Inches(i * 0.65)
        add_text_box(slide, before_x + Inches(0.4), y, card_w - Inches(0.8), Inches(0.4),
                    f"{icon}  {text}", font_size=13, font_color=Colors.MEDIUM_TEXT)
    
    # VS circle in center
    add_circle_icon(slide, Inches(6.6), Inches(4.0), Inches(0.4), Colors.SOFT_PURPLE, "VS", Colors.WHITE, 16)
    
    # AFTER card
    after_x = Inches(7.0)
    card = add_card(slide, after_x, before_y, card_w, card_h, Colors.VERY_LIGHT_MINT)
    card.line.fill.background()
    
    # After header
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, after_x, before_y, card_w, Inches(0.7), Colors.SOFT_MINT)
    add_text_box(slide, after_x + Inches(0.3), before_y + Inches(0.12), card_w - Inches(0.6), Inches(0.5),
                "😊 AFTER (With Time Management)", font_size=17, font_color=Colors.WHITE, bold=True)
    
    after_items = [
        ("✅", "Wakes up early, plans the day"),
        ("✅", "Limits phone to 30 minutes"),
        ("✅", "Finishes homework by evening"),
        ("✅", "Feels calm and confident"),
        ("✅", "Submits work on time"),
        ("✅", "Has time for sports & hobbies"),
    ]
    
    for i, (icon, text) in enumerate(after_items):
        y = before_y + Inches(1.0) + Inches(i * 0.65)
        add_text_box(slide, after_x + Inches(0.4), y, card_w - Inches(0.8), Inches(0.4),
                    f"{icon}  {text}", font_size=13, font_color=Colors.MEDIUM_TEXT)
    
    add_slide_number(slide, 9)


# ============================================================
# SLIDE 10: CONCLUSION
# ============================================================
def create_slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    add_top_accent_strip(slide)
    
    add_section_header(slide, "Conclusion", "Key takeaways from our presentation")
    
    # Summary diagram - circular flow (smaller radius to keep inside page)
    center_x = Inches(6.666)
    center_y = Inches(3.5)
    
    # Central circle
    add_circle_icon(slide, center_x, center_y, Inches(0.7), Colors.SOFT_PURPLE, "⭐", Colors.WHITE, 28)
    add_text_box(slide, center_x - Inches(0.8), center_y + Inches(0.75), Inches(1.6), Inches(0.3),
                "Success", font_size=13, font_color=Colors.DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Surrounding points - reduced radius so everything fits inside slide
    conclusions = [
        ("⏰", "Time Mgmt\nSaves Time", Colors.SOFT_BLUE, -60),
        ("😊", "Makes Life\nEasy", Colors.SOFT_MINT, 0),
        ("🌟", "Small Habits\n= Big Success", Colors.SOFT_PEACH, 60),
        ("🎯", "Stay Focused\n& Organized", Colors.SOFT_PINK, 120),
        ("📈", "Better Grades\n& Results", Colors.SOFT_TEAL, 180),
        ("🏆", "Achieve Your\nGoals", Colors.SOFT_LAVENDER, 240),
    ]
    
    radius = Inches(1.8)
    for icon, text, color, angle_deg in conclusions:
        angle_rad = math.radians(angle_deg - 90)
        px = center_x + radius * math.cos(angle_rad)
        py = center_y + radius * math.sin(angle_rad)
        
        # Outer circle
        add_circle_icon(slide, px, py, Inches(0.45), color, icon, Colors.WHITE, 20)
        
        # Label
        add_text_box(slide, px - Inches(0.65), py + Inches(0.5), Inches(1.3), Inches(0.55),
                    text, font_size=10, font_color=Colors.DARK_TEXT, alignment=PP_ALIGN.CENTER, bold=True)
    
    # Bottom quote - positioned to not overlap diagram
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.65), Colors.VERY_LIGHT_PURPLE)
    add_text_box(slide, Inches(1.8), Inches(6.65), Inches(10), Inches(0.5),
                "💡 \"The key is not to prioritize what's on your schedule, but to schedule your priorities.\"",
                font_size=13, font_color=Colors.ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
    
    add_slide_number(slide, 10)


# ============================================================
# SLIDE 11: THANK YOU
# ============================================================
def create_slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, Colors.WHITE)
    
    # Top accent
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08), Colors.SOFT_BLUE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.08), SLIDE_WIDTH, Inches(0.04), Colors.SOFT_LAVENDER)
    
    # Decorative circles - kept fully inside page, smaller
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), Inches(0.5), Inches(1.8), Inches(1.8), Colors.VERY_LIGHT_BLUE)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(11.0), Inches(0.5), Inches(1.8), Inches(1.8), Colors.VERY_LIGHT_PURPLE)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(0.5), Inches(5.2), Inches(1.8), Inches(1.8), Colors.VERY_LIGHT_MINT)
    add_shape(slide, MSO_SHAPE.OVAL, Inches(11.0), Inches(5.2), Inches(1.8), Inches(1.8), Colors.VERY_LIGHT_PINK)
    
    # Decorative dots
    add_decorative_dots(slide, Inches(2.5), Inches(1.5), 3, 3, Inches(0.3), Colors.SOFT_LAVENDER)
    add_decorative_dots(slide, Inches(10), Inches(5.0), 3, 3, Inches(0.3), Colors.SOFT_SKY)
    
    # Thank You text
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11.3), Inches(1.2),
                "Thank You!", font_size=60, font_color=Colors.DARK_TEXT, bold=True, alignment=PP_ALIGN.CENTER)
    
    # Divider
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.3), Inches(2.3), Inches(0.05), Colors.SOFT_BLUE)
    
    # Icons row
    icons_data = [
        ("⏰", Colors.SOFT_BLUE),
        ("📅", Colors.SOFT_PURPLE),
        ("🎯", Colors.SOFT_PINK),
        ("📊", Colors.SOFT_TEAL),
        ("🏆", Colors.SOFT_PEACH),
    ]
    
    for i, (icon, color) in enumerate(icons_data):
        ix = Inches(4.0) + Inches(i * 1.2)
        add_circle_icon(slide, ix, Inches(4.0), Inches(0.4), color, icon, Colors.WHITE, 20)
    
    # Any Questions?
    add_text_box(slide, Inches(1), Inches(4.8), Inches(11.3), Inches(0.6),
                "Any Questions?", font_size=28, font_color=Colors.ACCENT_PURPLE, alignment=PP_ALIGN.CENTER)
    
    # Team info
    add_text_box(slide, Inches(1), Inches(5.6), Inches(11.3), Inches(0.4),
                "Time Management  |  Work Smart, Save Time", font_size=14, 
                font_color=Colors.LIGHT_TEXT, alignment=PP_ALIGN.CENTER)
    
    # Bottom accent
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.38), SLIDE_WIDTH, Inches(0.12), Colors.SOFT_BLUE)
    
    add_slide_number(slide, 11)


# ============================================================
# GENERATE ALL SLIDES
# ============================================================
print("🎨 Creating Time Management Presentation...")
print("=" * 50)

create_slide_1(prs)
print("✅ Slide 1: Title Slide")

create_slide_2(prs)
print("✅ Slide 2: Introduction")

create_slide_3(prs)
print("✅ Slide 3: Importance of Time Management")

create_slide_4(prs)
print("✅ Slide 4: Common Problems")

create_slide_5(prs)
print("✅ Slide 5: Key Techniques")

create_slide_6(prs)
print("✅ Slide 6: Popular Methods")

create_slide_7(prs)
print("✅ Slide 7: Tools for Time Management")

create_slide_8(prs)
print("✅ Slide 8: Tips for Students")

create_slide_9(prs)
print("✅ Slide 9: Real-Life Example")

create_slide_10(prs)
print("✅ Slide 10: Conclusion")

create_slide_11(prs)
print("✅ Slide 11: Thank You")

# Save
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "Time_Management_Presentation_v4.pptx")
prs.save(output_path)
print("=" * 50)
print(f"🎉 Presentation saved to: {output_path}")
print(f"📊 Total Slides: {len(prs.slides)}")
print("✨ Done! Open the file in PowerPoint to view your presentation.")
